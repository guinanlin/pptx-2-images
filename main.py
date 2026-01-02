import os
import shutil
import subprocess
import tempfile
import uuid
import re
import time
from pathlib import Path
from typing import List, Set
from zipfile import ZipFile
from collections import defaultdict
from datetime import datetime, timedelta

from fastapi import FastAPI, UploadFile, File, HTTPException, Request, status
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from starlette.background import BackgroundTasks
from starlette.middleware.base import BaseHTTPMiddleware
from starlette.types import ASGIApp
from pptx import Presentation

app = FastAPI(
    title="PPTX to JPEG Converter Service",
    description="API to convert PPTX documents to JPEG images, one image per slide.",
    version="1.0.0"
)

# Create static files directory for serving images
STATIC_DIR = Path("/app/static")

# Mount static files
app.mount("/static", StaticFiles(directory=str(STATIC_DIR), html=True), name="static")

# --- Security Configuration ---

# 可疑路径黑名单（常见的扫描路径）
SUSPICIOUS_PATHS: Set[str] = {
    # 配置文件扫描
    "/.env", "/.env.bak", "/.env.save", "/.env.local", "/.env.production",
    "/backend/.env", "/admin/.env", "/api/.env",
    "/config.php", "/config.php.bak", "/config.js", "/aws-config.js", "/aws.config.js",
    "/wp-config.php", "/wp-config.php.old",
    # Git 相关
    "/.git/config", "/.git/HEAD", "/.git/",
    # WordPress 扫描
    "/wp-includes/", "/wp-admin/", "/xmlrpc.php",
    "/wp-login.php", "/wp-content/",
    # 其他常见扫描
    "/robots.txt", "/favicon.ico",  # 这些是正常的，但频繁请求可能是扫描
    "/.well-known/", "/.htaccess", "/.htpasswd",
    "/phpinfo.php", "/info.php", "/test.php",
    # 特定路径模式
    "/js/twint_ch.js", "/js/lkk_ch.js", "/css/support_parent.css",
}

# 速率限制配置
RATE_LIMIT_WINDOW = 60  # 时间窗口（秒）
RATE_LIMIT_MAX_REQUESTS = 100  # 每个IP在时间窗口内的最大请求数
RATE_LIMIT_STRICT_MAX = 20  # 严格限制：短时间内超过此数量直接拒绝

# 存储每个IP的请求时间戳
ip_request_times: defaultdict = defaultdict(list)
ip_blocked: Set[str] = set()

# --- Middleware for Security ---

class SecurityMiddleware(BaseHTTPMiddleware):
    """安全中间件：过滤可疑请求和速率限制"""
    
    async def dispatch(self, request: Request, call_next):
        client_ip = request.client.host if request.client else "unknown"
        path = request.url.path
        method = request.method
        
        # 检查是否被临时封禁
        if client_ip in ip_blocked:
            return JSONResponse(
                status_code=status.HTTP_429_TOO_MANY_REQUESTS,
                content={"detail": "Too many requests. Please try again later."}
            )
        
        # 速率限制检查
        current_time = time.time()
        # 清理过期的请求记录
        ip_request_times[client_ip] = [
            t for t in ip_request_times[client_ip] 
            if current_time - t < RATE_LIMIT_WINDOW
        ]
        
        # 检查速率限制
        if len(ip_request_times[client_ip]) >= RATE_LIMIT_STRICT_MAX:
            # 短时间内请求过多，临时封禁
            ip_blocked.add(client_ip)
            # 30分钟后自动解封
            def unblock():
                time.sleep(1800)  # 30分钟
                ip_blocked.discard(client_ip)
            import threading
            threading.Thread(target=unblock, daemon=True).start()
            return JSONResponse(
                status_code=status.HTTP_429_TOO_MANY_REQUESTS,
                content={"detail": "Rate limit exceeded. IP temporarily blocked."}
            )
        
        if len(ip_request_times[client_ip]) >= RATE_LIMIT_MAX_REQUESTS:
            return JSONResponse(
                status_code=status.HTTP_429_TOO_MANY_REQUESTS,
                content={"detail": "Too many requests. Please slow down."}
            )
        
        # 记录请求时间
        ip_request_times[client_ip].append(current_time)
        
        # 检查可疑路径
        is_suspicious = False
        for suspicious_path in SUSPICIOUS_PATHS:
            if path.startswith(suspicious_path) or suspicious_path in path:
                is_suspicious = True
                break
        
        # 对于可疑请求，直接返回404，不记录日志
        if is_suspicious:
            return JSONResponse(
                status_code=status.HTTP_404_NOT_FOUND,
                content={"detail": "Not Found"}
            )
        
        # 正常请求继续处理
        response = await call_next(request)
        return response

# --- Custom Logging Filter ---

class SecurityLogFilter:
    """自定义日志过滤器，过滤掉可疑请求的日志"""
    
    @staticmethod
    def should_log(path: str, status_code: int) -> bool:
        """判断是否应该记录日志"""
        # 对于404的可疑路径，不记录日志
        if status_code == 404:
            for suspicious_path in SUSPICIOUS_PATHS:
                if path.startswith(suspicious_path) or suspicious_path in path:
                    return False
        
        # 对于频繁的健康检查请求，可以降低日志级别（这里保留，但可以调整）
        if path in ["/", "/health"] and status_code == 200:
            # 可以在这里添加逻辑来减少健康检查的日志
            pass
        
        return True

# --- Custom Access Logging Middleware ---

class FilteredAccessLogMiddleware(BaseHTTPMiddleware):
    """过滤访问日志的中间件，减少日志噪音"""
    
    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        
        # 只记录非可疑请求的日志
        path = request.url.path
        status_code = response.status_code
        client_ip = request.client.host if request.client else "unknown"
        method = request.method
        
        # 判断是否应该记录日志
        should_log = SecurityLogFilter.should_log(path, status_code)
        
        if should_log:
            # 使用 print 输出日志（uvicorn 会捕获）
            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            print(f"{timestamp} - {client_ip} - {method} {path} - {status_code}")
        
        return response

# 注意：FastAPI 中间件是反向执行的（LIFO），最后添加的中间件会最先执行
# 所以先添加日志中间件，后添加安全中间件，这样安全中间件会先拦截可疑请求
# 可疑请求在安全中间件中被拦截后，不会到达日志中间件，因此不会记录日志
app.add_middleware(FilteredAccessLogMiddleware)
app.add_middleware(SecurityMiddleware)  # 这个会先执行，拦截可疑请求

# --- Utility Functions for Conversion ---

def sanitize_filename(filename: str) -> str:
    """Sanitize filename to avoid encoding issues with non-ASCII characters."""
    # Remove or replace problematic characters
    # Keep only alphanumeric, dots, hyphens, and underscores
    sanitized = re.sub(r'[^\w\-_\.]', '_', filename)
    # Ensure it's not empty and has a reasonable length
    if not sanitized or len(sanitized) > 100:
        # Generate a safe filename with UUID
        sanitized = f"file_{uuid.uuid4().hex[:8]}"
    return sanitized

def cleanup_path(path: Path):
    """Removes a file or directory after a response is sent."""
    if path.is_file():
        path.unlink(missing_ok=True)
    elif path.is_dir():
        shutil.rmtree(path, ignore_errors=True)

def cleanup_static_images(image_urls: List[str]):
    """Removes static images after a delay."""
    import time
    import threading
    
    def delayed_cleanup():
        time.sleep(3600)  # Wait 1 hour before cleanup
        for image_url in image_urls:
            image_filename = image_url.split('/')[-1]
            image_path = STATIC_DIR / image_filename
            if image_path.exists():
                image_path.unlink(missing_ok=True)
    
    # Run cleanup in a separate thread
    cleanup_thread = threading.Thread(target=delayed_cleanup)
    cleanup_thread.daemon = True
    cleanup_thread.start()

def extract_notes_from_pptx(pptx_path: Path) -> List[dict]:
    """Extract notes from PPTX file."""
    try:
        prs = Presentation(pptx_path)
        notes = []
        
        for i, slide in enumerate(prs.slides):
            # Get notes for this slide
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
            else:
                notes_text = ""
            
            notes.append({
                "slide_number": i + 1,
                "notes": notes_text.strip()
            })
        
        return notes
    except Exception as e:
        print(f"Warning: Could not extract notes from PPTX: {e}")
        return []

def convert_pptx_to_pdf(input_pptx_path: Path, output_dir: Path) -> Path:
    """Converts a PPTX file to PDF using soffice (LibreOffice)."""
    # LibreOffice's soffice command is used for conversion.
    # --headless: run without a graphical interface
    # --convert-to pdf: specify output format
    # --outdir: specify output directory
    output_pdf_path = output_dir / f"{input_pptx_path.stem}.pdf"
    
    try:
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(input_pptx_path)
        ]
        result = subprocess.run(cmd, capture_output=True, check=True, text=True, timeout=300, encoding='utf-8', errors='replace') # 5 min timeout
        print(f"Soffice stdout: {result.stdout}")
        print(f"Soffice stderr: {result.stderr}")

        if not output_pdf_path.exists():
            raise RuntimeError(f"PDF conversion failed: Output PDF not found. Stderr: {result.stderr}")
        return output_pdf_path
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"PPTX to PDF conversion failed: {e.stderr}")
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="PPTX to PDF conversion timed out.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX to PDF conversion error: {str(e)}")

def convert_pdf_to_jpegs(input_pdf_path: Path, output_dir: Path) -> List[Path]:
    """Converts a PDF file to JPEG images using ImageMagick."""
    # ImageMagick's convert command is used.
    # -density 150: sets the resolution (DPI) of the output images
    # -quality 80: sets the JPEG compression quality
    # output-%03d.jpg: specifies the output filename pattern for multiple pages (e.g., slide-001.jpg, slide-002.jpg)
    output_jpeg_pattern = output_dir / "slide-%03d.jpg"

    try:
        cmd = [
            "convert",
            "-density", "72",
            "-quality", "70",
            str(input_pdf_path),
            str(output_jpeg_pattern)
        ]
        result = subprocess.run(cmd, capture_output=True, check=True, text=True, timeout=300, encoding='utf-8', errors='replace') # 5 min timeout
        print(f"ImageMagick stdout: {result.stdout}")
        print(f"ImageMagick stderr: {result.stderr}")

        # Collect all generated JPEG files
        jpeg_files = sorted(list(output_dir.glob("slide-*.jpg")))
        if not jpeg_files:
            raise RuntimeError(f"PDF to JPEG conversion produced no images. Stderr: {result.stderr}")
        return jpeg_files
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"PDF to JPEG conversion failed: {e.stderr}")
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="PDF to JPEG conversion timed out.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to JPEG conversion error: {str(e)}")

# --- FastAPI Endpoints ---

@app.get("/")
async def root():
    """Health check endpoint."""
    return {"message": "PPTX to JPEG Converter Service is running", "status": "healthy"}

@app.get("/health")
async def health_check():
    """Health check endpoint for monitoring."""
    return {"status": "healthy", "service": "pptx-to-jpeg-converter"}

@app.get("/debug/static")
async def debug_static():
    """Debug endpoint to check static directory status."""
    try:
        # Check if static directory exists
        static_exists = STATIC_DIR.exists()
        static_is_dir = STATIC_DIR.is_dir() if static_exists else False
        
        # List files in static directory
        files = []
        if static_exists and static_is_dir:
            files = [f.name for f in STATIC_DIR.iterdir() if f.is_file()]
        
        return {
            "static_dir_path": str(STATIC_DIR),
            "static_dir_exists": static_exists,
            "static_dir_is_directory": static_is_dir,
            "files_count": len(files),
            "files": files[:10]  # Show first 10 files
        }
    except Exception as e:
        return {"error": str(e)}

@app.get("/debug/static/{filename}")
async def debug_static_file(filename: str):
    """Debug endpoint to check if a specific static file exists."""
    try:
        file_path = STATIC_DIR / filename
        exists = file_path.exists()
        is_file = file_path.is_file() if exists else False
        size = file_path.stat().st_size if exists and is_file else 0
        
        return {
            "filename": filename,
            "file_path": str(file_path),
            "exists": exists,
            "is_file": is_file,
            "size": size
        }
    except Exception as e:
        return {"error": str(e)}

@app.post("/convert/pptx-to-jpeg/")
async def convert_pptx_to_jpeg(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(..., description="The PPTX or PPT file to convert.")
):
    """
    Converts an uploaded PPTX/PPT file to JPEG images and returns URLs for each slide.
    Each slide is converted to a separate JPEG image and made available via static URLs.
    Images are automatically cleaned up after 1 hour.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file uploaded.")
    
    file_extension = Path(file.filename).suffix.lower()
    if file_extension not in (".pptx", ".ppt"):
        raise HTTPException(status_code=400, detail="Only PPTX or PPT files are supported.")

    # Create a temporary directory for processing this request
    # This ensures isolation and easier cleanup
    temp_dir = Path(tempfile.mkdtemp())
    
    # Sanitize the filename to avoid encoding issues
    safe_filename = sanitize_filename(file.filename)
    input_pptx_path = temp_dir / safe_filename

    # Ensure the temporary directory is cleaned up after the request
    background_tasks.add_task(cleanup_path, temp_dir)

    # Save the uploaded PPTX file to the temporary directory
    try:
        with open(input_pptx_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save uploaded file: {str(e)}")

    try:
        # Step 1: Extract notes from PPTX (before conversion)
        notes_data = extract_notes_from_pptx(input_pptx_path)
        
        # Step 2: PPTX to PDF
        pdf_path = convert_pptx_to_pdf(input_pptx_path, temp_dir)
        
        # Step 3: PDF to JPEG
        jpeg_files = convert_pdf_to_jpegs(pdf_path, temp_dir)

        # Copy images to static directory and generate URLs
        # Ensure static directory exists
        STATIC_DIR.mkdir(exist_ok=True, mode=0o755)
        
        # Generate a unique random ID for this conversion session
        session_id = uuid.uuid4().hex[:8]
        image_urls = []
        
        for i, jpeg_file in enumerate(jpeg_files):
            # Create a unique filename for each image using random ID + sequence
            image_filename = f"{session_id}_{i+1:03d}.jpg"
            static_image_path = STATIC_DIR / image_filename
            
            # Copy the image to static directory
            shutil.copy2(jpeg_file, static_image_path)
            
            # Set proper permissions for the image file
            static_image_path.chmod(0o644)
            
            # Debug: Print file info
            print(f"DEBUG: Copied {jpeg_file} to {static_image_path}")
            print(f"DEBUG: File exists: {static_image_path.exists()}")
            print(f"DEBUG: File size: {static_image_path.stat().st_size if static_image_path.exists() else 'N/A'}")
            
            # Generate URL for the image
            image_url = f"/static/{image_filename}"
            image_urls.append(image_url)
        
        # Create slides data with images and notes
        slides_data = []
        for i, (image_url, note_data) in enumerate(zip(image_urls, notes_data)):
            slides_data.append({
                "slide_number": i + 1,
                "image_url": image_url,
                "notes": note_data["notes"]
            })
        
        # Schedule cleanup of static images after 1 hour (3600 seconds)
        # Note: In production, you might want to use a proper task queue
        background_tasks.add_task(cleanup_static_images, image_urls)
        
        # Return JSON response with slides data
        return {
            "status": "success",
            "message": f"Successfully converted {len(jpeg_files)} slides",
            "slide_count": len(jpeg_files),
            "slides": slides_data,
            "original_filename": file.filename,
            "note": "Images will be automatically cleaned up after 1 hour"
        }
    except HTTPException:
        # Re-raise HTTPExceptions directly
        raise
    except Exception as e:
        # Catch any other unexpected errors during conversion
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred during conversion: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    import logging
    
    # 配置日志级别，减少噪音
    logging.getLogger("uvicorn.access").setLevel(logging.WARNING)
    
    uvicorn.run(
        app, 
        host="0.0.0.0", 
        port=8000,
        access_log=False  # 禁用默认访问日志，使用自定义中间件
    )